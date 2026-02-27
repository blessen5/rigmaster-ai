from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_table_slide(prs, title, headers, data):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Calculate table dimensions
    rows = len(data) + 1  # +1 for header
    cols = len(headers)
    
    # Add table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    col_widths = [Inches(1.8), Inches(1.2), Inches(0.8), Inches(2.5), Inches(1.0), Inches(1.7)]
    for i, width in enumerate(col_widths[:cols]):
        table.columns[i].width = width
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        
        # Format header text
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add data
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            
            # Format data text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "RigMaster AI Database Schema", "Complete MongoDB Collections Documentation")

# Slide 2: Database Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Database Overview"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Total Collections: 8"

p = tf.add_paragraph()
p.text = "1. users - User accounts and authentication"
p.level = 1

p = tf.add_paragraph()
p.text = "2. saved_builds - User PC build configurations"
p.level = 1

p = tf.add_paragraph()
p.text = "3. components - Hardware components catalog"
p.level = 1

p = tf.add_paragraph()
p.text = "4. ai_cache - AI recommendation cache"
p.level = 1

p = tf.add_paragraph()
p.text = "5. shopping_cache - Price lookup cache"
p.level = 1

p = tf.add_paragraph()
p.text = "6. price_alerts - User price tracking alerts"
p.level = 1

p = tf.add_paragraph()
p.text = "7. otps - One-time passwords for password reset"
p.level = 1

# Slide 3: users table
headers = ["Field Name", "Data Type", "Size", "Constraints", "Default", "Description"]
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique user ID"],
    ["email", "VARCHAR", "255", "NOT NULL, UNIQUE", "NULL", "User's email"],
    ["username", "VARCHAR", "50", "NOT NULL, UNIQUE", "NULL", "Username"],
    ["password", "VARCHAR", "255", "NOT NULL", "NULL", "Hashed password"],
    ["is_admin", "BOOLEAN", "1", "NOT NULL", "FALSE", "Admin flag"],
    ["is_active", "BOOLEAN", "1", "NOT NULL", "TRUE", "Active status"],
    ["created_at", "DATETIME", "8 bytes", "NOT NULL", "NOW()", "Creation timestamp"],
]
add_table_slide(prs, "TABLE: users", headers, data)

# Slide 4: saved_builds table
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique build ID"],
    ["user_id", "VARCHAR", "24", "NOT NULL, FOREIGN KEY → users._id", "NULL", "Build owner"],
    ["name", "VARCHAR", "100", "NOT NULL", "'Custom Rig'", "Build name"],
    ["cpu_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "CPU reference"],
    ["gpu_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "GPU reference"],
    ["motherboard_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "Motherboard ref"],
    ["ram_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "RAM reference"],
    ["storage_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "Storage ref"],
    ["psu_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "PSU reference"],
    ["case_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "Case reference"],
    ["cooler_id", "ObjectId", "12 bytes", "NULL, FOREIGN KEY → components._id", "NULL", "Cooler ref"],
    ["is_public", "BOOLEAN", "1", "NOT NULL", "FALSE", "Public visibility"],
    ["created_at", "DATETIME", "8 bytes", "NOT NULL", "NOW()", "Creation time"],
]
add_table_slide(prs, "TABLE: saved_builds", headers, data)

# Slide 5: components table (common fields)
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique component ID"],
    ["name", "VARCHAR", "255", "NOT NULL", "NULL", "Component name"],
    ["category", "ENUM", "-", "NOT NULL", "NULL", "Type: cpu, gpu, etc."],
    ["brand", "VARCHAR", "100", "NULL", "NULL", "Manufacturer"],
    ["status", "VARCHAR", "20", "NULL", "'Active'", "Availability"],
    ["price", "DECIMAL", "(10,2)", "NULL", "NULL", "Price in USD"],
    ["msrp", "DECIMAL", "(10,2)", "NULL", "NULL", "MSRP in USD"],
    ["retailer", "VARCHAR", "100", "NULL", "NULL", "Retailer name"],
    ["product_url", "TEXT", "2048", "NULL", "NULL", "Product URL"],
    ["in_stock", "BOOLEAN", "1", "NULL", "TRUE", "Stock status"],
]
add_table_slide(prs, "TABLE: components (Common Fields)", headers, data)

# Slide 6: components - CPU specific fields
data = [
    ["socket", "VARCHAR", "50", "NULL", "NULL", "CPU socket type"],
    ["cores", "INT", "4", "NULL", "NULL", "Number of cores"],
    ["core_count", "INT", "4", "NULL", "NULL", "Core count"],
    ["boost_clock", "VARCHAR", "50", "NULL", "NULL", "Boost clock speed"],
    ["core_clock", "VARCHAR", "50", "NULL", "NULL", "Base clock speed"],
    ["tdp", "INT", "4", "NULL", "NULL", "TDP in Watts"],
    ["microarchitecture", "VARCHAR", "100", "NULL", "NULL", "CPU architecture"],
    ["graphics", "VARCHAR", "100", "NULL", "NULL", "Integrated GPU"],
]
add_table_slide(prs, "TABLE: components (CPU Fields)", headers, data)

# Slide 7: components - GPU, Motherboard, RAM fields
data = [
    ["vram", "VARCHAR", "20", "NULL", "NULL", "GPU: Video memory"],
    ["socket", "VARCHAR", "50", "NULL", "NULL", "Mobo: CPU socket"],
    ["form_factor", "VARCHAR", "20", "NULL", "NULL", "Mobo: Size format"],
    ["memory_type", "VARCHAR", "10", "NULL", "NULL", "Mobo: RAM type"],
    ["capacity", "VARCHAR", "20", "NULL", "NULL", "RAM: Memory size"],
    ["speed", "VARCHAR", "20", "NULL", "NULL", "RAM: Speed (MHz)"],
]
add_table_slide(prs, "TABLE: components (GPU/Mobo/RAM)", headers, data)

# Slide 8: components - Storage, PSU fields
data = [
    ["capacity", "VARCHAR", "20", "NULL", "NULL", "Storage: Capacity"],
    ["type", "VARCHAR", "20", "NULL", "NULL", "Storage: Type"],
    ["wattage", "INT", "4", "NULL", "NULL", "PSU: Power output"],
]
add_table_slide(prs, "TABLE: components (Storage/PSU)", headers, data)

# Slide 9: ai_cache table
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique cache ID"],
    ["cache_key", "VARCHAR", "255", "NOT NULL, UNIQUE", "NULL", "Cache identifier"],
    ["build", "JSON", "16MB max", "NULL", "NULL", "Build recommendation"],
    ["total_estimated_cost", "DECIMAL", "(10,2)", "NULL", "NULL", "Total cost"],
    ["explanation", "TEXT", "65535", "NULL", "NULL", "AI explanation"],
    ["prediction", "JSON", "16MB max", "NULL", "NULL", "AI prediction data"],
    ["created_at", "DATETIME", "8 bytes", "NOT NULL", "NOW()", "Cache timestamp"],
]
add_table_slide(prs, "TABLE: ai_cache", headers, data)

# Slide 10: shopping_cache table
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique cache ID"],
    ["query", "VARCHAR", "255", "NOT NULL, UNIQUE", "NULL", "Search query"],
    ["listings", "JSON", "16MB max", "NOT NULL", "[]", "Product listings"],
    ["expires_at", "DATETIME", "8 bytes", "NOT NULL", "NOW()+24h", "Expiration time"],
]
add_table_slide(prs, "TABLE: shopping_cache", headers, data)

# Slide 11: price_alerts table
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique alert ID"],
    ["user_id", "ObjectId", "12 bytes", "NOT NULL, FOREIGN KEY → users._id", "NULL", "User reference"],
    ["component_id", "ObjectId", "12 bytes", "NOT NULL, FOREIGN KEY → components._id", "NULL", "Component ref"],
    ["target_price", "DECIMAL", "(10,2)", "NOT NULL, CHECK > 0", "NULL", "Price threshold"],
    ["triggered", "BOOLEAN", "1", "NOT NULL", "FALSE", "Alert status"],
    ["created_at", "DATETIME", "8 bytes", "NOT NULL", "NOW()", "Creation time"],
]
add_table_slide(prs, "TABLE: price_alerts", headers, data)

# Slide 12: otps table
data = [
    ["_id", "ObjectId", "12 bytes", "PRIMARY KEY, AUTO_INCREMENT, NOT NULL, UNIQUE", "Auto", "Unique OTP ID"],
    ["email", "VARCHAR", "255", "NOT NULL, UNIQUE", "NULL", "User's email"],
    ["otp", "CHAR", "6", "NOT NULL, CHECK: [0-9]{6}", "NULL", "6-digit OTP"],
    ["expires_at", "DATETIME", "8 bytes", "NOT NULL", "NOW()+10m", "Expiration time"],
]
add_table_slide(prs, "TABLE: otps", headers, data)

# Slide 13: Database Indexes
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Database Indexes"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "users"

p = tf.add_paragraph()
p.text = "UNIQUE INDEX on email"
p.level = 1

p = tf.add_paragraph()
p.text = "UNIQUE INDEX on username"
p.level = 1

p = tf.add_paragraph()
p.text = "saved_builds"
p.level = 0

p = tf.add_paragraph()
p.text = "INDEX on user_id"
p.level = 1

p = tf.add_paragraph()
p.text = "COMPOSITE INDEX on (user_id, created_at DESC)"
p.level = 1

p = tf.add_paragraph()
p.text = "components"
p.level = 0

p = tf.add_paragraph()
p.text = "INDEX on category"
p.level = 1

p = tf.add_paragraph()
p.text = "UNIQUE INDEX on (category, name)"
p.level = 1

p = tf.add_paragraph()
p.text = "ai_cache / shopping_cache"
p.level = 0

p = tf.add_paragraph()
p.text = "UNIQUE INDEX on cache_key / query"
p.level = 1

p = tf.add_paragraph()
p.text = "TTL INDEX on expires_at"
p.level = 1

# Slide 14: Global Constraints
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Global Application Constraints"
content = slide.placeholders[1]
tf = content.text_frame

constraints = [
    ("Max Request Body Size", "16 MB"),
    ("Password Min Length", "6 characters"),
    ("Password Max Length", "255 characters"),
    ("Username Max Length", "50 characters"),
    ("Email Max Length", "255 characters"),
    ("Build Name Max Length", "100 characters"),
    ("OTP Length", "Exactly 6 digits"),
    ("OTP Expiry", "10 minutes"),
    ("Shopping Cache TTL", "24 hours"),
    ("AI Cache TTL", "7 days"),
    ("Budget Range", "$400 - $6,000"),
    ("Build Comparison Limit", "2-3 builds"),
]

for constraint, value in constraints:
    p = tf.add_paragraph()
    p.text = f"{constraint}: {value}"
    p.font.size = Pt(14)

# Save presentation
output_file = "RigMaster_Database_Schema.pptx"
prs.save(output_file)
print(f"[SUCCESS] PowerPoint presentation created successfully: {output_file}")
print(f"Total slides: {len(prs.slides)}")
print(f"File location: {output_file}")
